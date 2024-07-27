import re 
import requests 
from bs4 import BeautifulSoup 
from pptx import Presentation 
from pptx.util import Pt, Inches 
from pptx.enum.text import PP_ALIGN 
from pptx.dml.color import RGBColor 

def create_presentation():
    return Presentation()

def add_title_slide(prs, title_text, background = RGBColor(31,56,100)):
    if not title_text:
        return
    title_slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.add_textbox(left=Inches(1), top=Inches(0.5), width=Inches(14), height=Inches(4))
    title.text = title_text
    title.text_frame.text = title_text
    index = 0
    title.text_frame.margin_top = Inches(2.75)

    for paragraph in title.text_frame.paragraphs:
        if index == 0:    
            paragraph.alignment = PP_ALIGN.CENTER
            title.text_frame.word_wrap = True
            if paragraph.runs:
                title_format = paragraph.runs[0].font
            title_format.bold = True
            title_format.name = 'Century Gothic'
            title_format.size = Pt(80)
            title_format.color.rgb = RGBColor(255, 255, 255)
        else:
            paragraph.alignment = PP_ALIGN.CENTER
            title.text_frame.word_wrap = True
            if paragraph.runs:
                title_format = paragraph.runs[0].font
            title_format.bold = False
            title_format.name = 'Century Gothic'
            title_format.size = Pt(60)
            title_format.color.rgb = RGBColor(255, 255, 255)
            title.text_frame.margin_top = Inches(1)

        index = index + 1
    
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = background

def add_content_slide(prs, content_text = '', max_chars_per_line=210, background = RGBColor(31,56,100)):
    if not content_text:
        return
    
    content_slide_layout = prs.slide_layouts[6]
    content_text.replace('\n', '\n')
    words = content_text.split()
    lines = []
    current_line = ''

    for word in words:
        if word == '<br>':
            lines.append(current_line)
            current_line = ' '
        else:
            if len(current_line) + len(word) + 1 <= max_chars_per_line:
                current_line += word + ' '
            else:
                lines.append(current_line.strip())
                current_line = word + ' '

    if current_line:
        lines.append(current_line.strip())

    for line in lines:
        if not line:
            line = " "

        slide = prs.slides.add_slide(content_slide_layout)
        content = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(0.5), width=Inches(15), height=Inches(8))
        content.text = line
        content.text_frame.text = line        

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = background

        if "AS:" in line:
            adjust_text_box_size(content, True, len(line) > 86)
        else:
            adjust_text_box_size(content, False, False)

def adjust_text_box_size(shape, isCenter, isLarge):
    text_frame = shape.text_frame

    for paragraph in shape.text_frame.paragraphs:
        if not paragraph.runs:
            return
        font_format = paragraph.runs[0].font
        font_format.color.rgb = RGBColor(255,255,255)
        font_format.bold = isCenter
        font_format.name = 'Century Gothic'
        if isCenter:
            font_format.size = Pt(65)
        else:    
            font_format.size = Pt(60)
        paragraph.alignment = PP_ALIGN.LEFT

    if isCenter:
        if isLarge:
            text_frame.margin_top = Inches(1.25) 
        else:
            text_frame.margin_top = Inches(3) 


    text_frame.word_wrap = True
    text_frame.margin_right = 0  
    text_frame.margin_left = 0  


def save_presentation(prs, file_path):
    prs.save(file_path)

def get_parts(html, htmlId, get_text, get_sub = False, stepTitlePosition = 0, stepSubPosition = 0):
    sub = ''
    text = ''
    title = ''
    title_founded = ''
    div = html.find(id=f"{htmlId}")

    if not div:
        return {'title': title, 'text': text, 'sub': sub}
    if div.find('b'):
        title_founded = div.find_all('b')[0 + stepTitlePosition].getText()
    elif div.find('strong'):
        title_founded = div.find_all('strong')[0 + stepTitlePosition].getText()
    else:
        ImportError()
        
    if get_sub:
        sub_founded = div.find_all('p')[1 + stepSubPosition].getText()
        
        if sub_founded and not (sub_founded.startswith('http')):
            sub = sub_founded
        elif div.find_all('p')[2 + stepSubPosition].getText():
            sub = div.find_all('p')[2 + stepSubPosition].getText()

    if title_founded:
        title = title_founded 

    if get_text:
        paragraphs = div.find_all('p')[1:]
        text = '\n'.join([p.getText() for p in paragraphs])
        text = re.sub(r'(?<=[a-zA-Zá-üÁ-Ü\s])\d+|\d+(?=[a-zA-Zá-üÁ-Ü\s])', '', text, flags=re.UNICODE)

    return {'title': title, 'text': text, 'sub': sub}

def get_txt_file(path):
    try:
        with open(path, 'r', encoding='utf-8') as file:
            text = file.read()
        return text
    except FileNotFoundError:
        print(f"O arquivo '{path}' não foi encontrado.")
    except Exception as e:
        print(f"Ocorreu um erro ao ler o arquivo: {e}")

def main():
    ###### ALTERE AQUI
    ###### ALTERE AQUI
    # Após alteração precione "CTRL + S"
    
    url = 'https://liturgia.cancaonova.com/pb/liturgia/11o-domingo-tempo-comum-domingo/?sDia=16&sMes=06&sAno=2024' # ver site da canção nova
    oracao = get_txt_file('oracoes_eucaristicas/eucaristica_4_salvador.txt') # ver pasta de arquivos "oracoes_eucaristicas"
    comGloria = True
    comCreio = True
    # primeiraEucaristia = get_txt_file('outras_oracoes/primeira_eucaristia_renovacao.txt')
    primeiraEucaristia = False

    # Após alteração precione "CTRL + S"
    ###### ALTERE AQUI
    ###### ALTERE AQUI
    
    response = requests.get(url)
    html_content = response.content
    html = BeautifulSoup(html_content, 'html.parser')
    title = html.select_one('h1.entry-title').getText().strip()

    l1 = get_parts(html, 'liturgia-1', False, True, 1, 1)
    s = get_parts(html, 'liturgia-2', False, True, 1, 1)
    l2 = get_parts(html, 'liturgia-3', False, True, 1, 1)
    e = get_parts(html, 'liturgia-4', False, False)

    pai_nosso = get_txt_file('outras_oracoes/pai_nosso.txt')
    creio = get_txt_file('outras_oracoes/creio.txt')

    presentation = create_presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)

    add_title_slide(presentation, title)
    add_title_slide(presentation, "CANTO DE ENTRADA")
    add_title_slide(presentation, " ")
    add_title_slide(presentation, "ATO PENITENCIAL")
    add_title_slide(presentation, " ")

    if comGloria:
        add_title_slide(presentation, "GLÓRIA")
        add_title_slide(presentation, " ")

    add_title_slide(presentation, l1['title'] + '\n \n' + l1['sub'])
    add_content_slide(presentation, l1['text'])
    add_title_slide(presentation, " ")
    add_title_slide(presentation, s['title'])
    add_title_slide(presentation, s['sub'])
    add_content_slide(presentation, s['text'])
    add_title_slide(presentation, " ")

    if l2['title']:
        add_title_slide(presentation, l2['title'] + '\n \n' + l2['sub'])
        add_content_slide(presentation, l2['text'])

    add_title_slide(presentation, " ")
    add_title_slide(presentation, "ACLAMAÇÃO AO EVANGELHO")
    add_title_slide(presentation, e['title'])
    add_content_slide(presentation, e['text'])
    add_title_slide(presentation, " ")
    add_title_slide(presentation, "HOMILIA")

    if comCreio:
        add_content_slide(presentation, creio)
        add_title_slide(presentation, " ")
        
    if primeiraEucaristia:
        add_content_slide(presentation, primeiraEucaristia)
    add_title_slide(presentation, " ")

    add_title_slide(presentation, "PRECES DA COMUNIDADE" + "/n /n" + "Cristo, rei da glória, ouvi-nos!")
    add_title_slide(presentation, " ")

    add_title_slide(presentation, "OFERTÓRIO")
    add_title_slide(presentation, " ")
    add_title_slide(presentation, "ORAÇÃO EUCARISTICA")
    add_content_slide(presentation, oracao)
    add_title_slide(presentation, " ")
    add_content_slide(presentation, pai_nosso)
    add_title_slide(presentation, " ")
    add_title_slide(presentation, "CORDEIRO")
    add_title_slide(presentation, " ")
    add_title_slide(presentation, "COMUNHÃO")
    add_title_slide(presentation, " ")
    add_title_slide(presentation, " ")

    save_presentation(presentation, "Missa textos.pptx")

def mainSabadoAleluia():
    ###### ALTERE AQUI
    ###### ALTERE AQUI
    # Após alteração precione "CTRL + S"
    url = 'https://liturgia.cancaonova.com/pb/liturgia/sabado-santo-vigilia-pascal-5/?sDia=30&sMes=03&sAno=2024' # ver site da canção nova
    oracao = get_txt_file('oracoes_eucaristicas/eucaristica_1_pascoal_salvador.txt') # ver pasta de arquivos "oracoes_eucaristicas"
    # Após alteração precione "CTRL + S"
    ###### ALTERE AQUI
    ###### ALTERE AQUI
    
    response = requests.get(url)
    html_content = response.content
    html = BeautifulSoup(html_content, 'html.parser')
    title = html.select_one('h1.entry-title').getText().strip()

    e = get_parts(html, 'liturgia-4', False, False)
    pai_nosso = get_txt_file('outras_oracoes/pai_nosso.txt')
    creio = get_txt_file('outras_oracoes/creio.txt')
    renovacaoBatismo = get_txt_file('outras_oracoes/renovacao_batismo.txt')
    ladainhaSantos = get_txt_file('outras_oracoes/ladainha_santos.txt')
    proclamacaoPascoa = get_txt_file('outras_oracoes/proclamacao_pascoa.txt')


    presentation = create_presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)

    add_title_slide(presentation, title)
    add_title_slide(presentation, "BENÇÃO DO FOGO", RGBColor(0,0,0))
    add_title_slide(presentation, " ", RGBColor(0,0,0))
    add_title_slide(presentation, "PROCLAMAÇÃO DA PÁSCOA", RGBColor(0,0,0))
    add_content_slide(presentation, proclamacaoPascoa, 210, RGBColor(0,0,0))
    add_title_slide(presentation, " ", RGBColor(0,0,0))
    add_title_slide(presentation, 'Primeira Leitura \n (Gn 1,1– 2,2)' + '\n \n' + 'Leitura do Livro do Gênesis:', RGBColor(0,0,0))
    add_title_slide(presentation, 'Salmo (103)' + '\n \n' + '— Enviai o vosso Espírito, Senhor, e da terra toda a face renovai.', RGBColor(0,0,0))
    add_title_slide(presentation, " ", RGBColor(0,0,0))
    add_title_slide(presentation, 'Segunda Leitura \n (Gn 22,1-2.9a.10-13.15-18)' + '\n \n' + 'Leitura do Livro do Gênesis.', RGBColor(0,0,0))
    add_title_slide(presentation, 'Salmo (15)' + '\n \n' + '— Guardai-me, ó Deus, porque em vós me refugio!', RGBColor(0,0,0))
    add_title_slide(presentation, " ", RGBColor(0,0,0))
    add_title_slide(presentation, 'Terceira Leitura \n (Êx 14,15 – 15,1)' + '\n \n' + 'Leitura do Livro do Êxodo.', RGBColor(0,0,0))
    add_title_slide(presentation, 'Salmo (Êx 15,1-6.17-18)' + '\n \n' + '— Cantemos ao Senhor que fez brilhar a sua glória!', RGBColor(0,0,0))
    add_title_slide(presentation, " ", RGBColor(0,0,0))
    add_title_slide(presentation, 'Quarta Leitura \n (Is 54, 5-14)' + '\n \n' + 'Leitura da Livro do Profeta Isaías.', RGBColor(0,0,0))
    add_title_slide(presentation, 'Salmo (Sl 29)' + '\n \n' + '— Eu vos exalto, ó Senhor, porque vós me livrastes!', RGBColor(0,0,0))
    add_title_slide(presentation, " ", RGBColor(0,0,0))
    add_title_slide(presentation, 'Quinta Leitura \n (Is 55,1-11)' + '\n \n' + 'Leitura da Livro do Profeta Isaías.', RGBColor(0,0,0))
    add_title_slide(presentation, 'Salmo (Is 12,2-6)' + '\n \n' + '— Com alegria bebereis do manancial da salvação.', RGBColor(0,0,0))
    add_title_slide(presentation, " ", RGBColor(0,0,0))
    add_title_slide(presentation, 'Sexta Leitura \n (Br 3,9-15.32–4,4)' + '\n \n' + 'Leitura do Livro do Profeta Baruc:', RGBColor(0,0,0))
    add_title_slide(presentation, 'Salmo (Sl 18)' + '\n \n' + '— Senhor, tens palavras de vida eterna.', RGBColor(0,0,0))
    add_title_slide(presentation, " ", RGBColor(0,0,0))
    add_title_slide(presentation, 'Sétima Leitura \n (Ez 36,16-17a.18-28)' + '\n \n' + 'Leitura da Profecia de Ezequiel:', RGBColor(0,0,0))
    add_title_slide(presentation, 'Salmo (Sl 41)' + '\n \n' + '— A minha alma tem sede de Deus.', RGBColor(0,0,0))
    add_title_slide(presentation, " ", RGBColor(0,0,0))

    add_title_slide(presentation, "GLÓRIA", RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))
    add_title_slide(presentation, 'Carta \n (Ro 6,3-11)' + '\n \n' + 'Leitura da carta de São Paulo aos Romanos:', RGBColor(118, 92, 75))
    add_title_slide(presentation, 'Salmo (117)' + '\n \n' + 'Aleluia, aleluia, aleluia.', RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))
    add_title_slide(presentation, e['title'], RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))
    add_title_slide(presentation, "HOMILIA", RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))
    add_content_slide(presentation, ladainhaSantos, 210, RGBColor(118, 92, 75))

    add_title_slide(presentation, " ", RGBColor(118, 92, 75))
    add_content_slide(presentation, creio, 210, RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))
    add_title_slide(presentation, "RENOVAÇÃO DAS PROMESSAS DO BATISMO", RGBColor(118, 92, 75))
    add_content_slide(presentation, renovacaoBatismo, 210, RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))

    add_title_slide(presentation, "PRECES DA COMUNIDADE", RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))

    add_title_slide(presentation, "OFERTÓRIO" + "/n /n" + "— Vinde, Senhor, com vossa luz.", RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))
    add_title_slide(presentation, "ORAÇÃO EUCARISTICA I", RGBColor(118, 92, 75))
    add_content_slide(presentation, oracao, 210, RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))
    add_content_slide(presentation, pai_nosso, 210, RGBColor(118, 92, 75)) 
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))
    add_title_slide(presentation, "CORDEIRO", RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))
    add_title_slide(presentation, "COMUNHÃO", RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))
    add_title_slide(presentation, " ", RGBColor(118, 92, 75))

    save_presentation(presentation, "Missa textos.pptx")

def mainCrisma():
    ###### ALTERE AQUI
    
    url = 'https://liturgia.cancaonova.com/pb/liturgia/17a-semana-tempo-comum-domingo/?sDia=28&sMes=07&sAno=2024' # ver site da canção nova
    oracao = get_txt_file('oracoes_eucaristicas/eucaristica_2_anunciamos.txt') # ver pasta de arquivos "oracoes_eucaristicas"
    comGloria = True
    comCreio = True
    renovacaoPromessasBatismo = get_txt_file('outras_oracoes/renovacao_batismo_crisma.txt')

    ## Musicas
    acenderVelas = get_txt_file('musicas/acendimento_das_velas.txt')
    entrada = get_txt_file('musicas/entrada.txt')
    atoPenitencial = get_txt_file('musicas/ato_penitencial.txt')
    gloria = get_txt_file('musicas/gloria.txt')
    ofertorio = get_txt_file('musicas/ofertorio.txt')
    comunhao = get_txt_file('musicas/comunhao.txt')
    uncaoCrisma = get_txt_file('musicas/uncao_do_crisma.txt')

    ###### ALTERE AQUI
    
    response = requests.get(url)
    html_content = response.content
    html = BeautifulSoup(html_content, 'html.parser')
    title = html.select_one('h1.entry-title').getText().strip()

    l1 = get_parts(html, 'liturgia-1', False, True, 1, 1)
    s = get_parts(html, 'liturgia-2', False, True, 0, 1)
    l2 = get_parts(html, 'liturgia-3', False, True, 1, 1)
    e = get_parts(html, 'liturgia-4', False, False)

    print(l1)
    print(s)
    print(l2)
    print(e)

    pai_nosso = get_txt_file('outras_oracoes/pai_nosso.txt')
    creio = get_txt_file('outras_oracoes/creio.txt')

    presentation = create_presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)

    add_title_slide(presentation, title)
    add_content_slide(presentation, entrada)
    add_title_slide(presentation, " ")
    add_title_slide(presentation, "ATO PENITENCIAL")
    add_content_slide(presentation, atoPenitencial)
    add_title_slide(presentation, " ")

    if comGloria:
        add_title_slide(presentation, "GLÓRIA")
        add_content_slide(presentation, gloria)
        add_title_slide(presentation, " ")

    add_title_slide(presentation, l1['title'] + '\n \n' + l1['sub'])
    add_content_slide(presentation, l1['text'])
    add_title_slide(presentation, " ")
    add_title_slide(presentation, s['title'])
    add_title_slide(presentation, s['sub'])
    # add_content_slide(presentation, s['text'])
    add_title_slide(presentation, " ")

    if l2['title']:
        add_title_slide(presentation, l2['title'] + '\n \n' + l2['sub'])
        add_content_slide(presentation, l2['text'])

    add_title_slide(presentation, " ")
    add_title_slide(presentation, "ACLAMAÇÃO AO EVANGELHO")
    add_title_slide(presentation, e['title'])
    add_content_slide(presentation, e['text'])
    add_title_slide(presentation, " ")
    add_title_slide(presentation, "HOMILIA")
    add_title_slide(presentation, " ")

    add_content_slide(presentation, acenderVelas)
    add_title_slide(presentation, " ")

    add_content_slide(presentation, renovacaoPromessasBatismo)
    add_title_slide(presentation, " ")

    if comCreio:
        add_content_slide(presentation, creio)
        add_title_slide(presentation, " ")

    add_title_slide(presentation, "UNÇÃO DO CRISMA")
    add_content_slide(presentation, uncaoCrisma)        
    add_title_slide(presentation, " ")


    add_title_slide(presentation, "PRECES DA COMUNIDADE" + "/n /n" + "Saciai, Senhor, o vosso povo!")
    add_title_slide(presentation, " ")

    add_title_slide(presentation, "OFERTÓRIO")
    add_content_slide(presentation, ofertorio)

    add_title_slide(presentation, " ")
    add_title_slide(presentation, "ORAÇÃO EUCARISTICA II")
    add_content_slide(presentation, oracao)
    add_title_slide(presentation, " ")
    add_content_slide(presentation, pai_nosso)
    add_title_slide(presentation, " ")
    add_title_slide(presentation, "CORDEIRO")
    add_title_slide(presentation, " ")
    add_title_slide(presentation, "COMUNHÃO")
    add_content_slide(presentation, comunhao)
    add_title_slide(presentation, " ")
    add_title_slide(presentation, " ")

    save_presentation(presentation, "Missa textos.pptx")

if __name__ == "__main__":
    mainCrisma()
